"""Document extraction: PDFs (pypdf) and PPTX (python-pptx) -> plain text.

Preserves [page N] markers so downstream citation logic can resolve them.
"""

from __future__ import annotations

import re
from pathlib import Path

from pypdf import PdfReader
from pptx import Presentation

_CONTROL_CHARS = re.compile(r"[\x00-\x08\x0b\x0c\x0e-\x1f]")


def _sanitize(text: str) -> str:
    return _CONTROL_CHARS.sub("", text)


def extract_pdf(path: Path) -> str:
    reader = PdfReader(str(path))
    parts: list[str] = []
    for i, page in enumerate(reader.pages, start=1):
        try:
            text = page.extract_text() or ""
        except Exception:
            text = ""
        parts.append(f"[page {i}]\n{_sanitize(text).strip()}")
    return "\n\n".join(parts)


def extract_pptx(path: Path) -> str:
    prs = Presentation(str(path))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in paragraph.runs if run.text)
                    if line.strip():
                        slide_texts.append(line.strip())
        parts.append(f"[page {i}]\n" + "\n".join(slide_texts))
    return "\n\n".join(parts)


def extract(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return extract_pdf(path)
    if suffix == ".pptx":
        return extract_pptx(path)
    raise ValueError(f"Unsupported file type: {path}")
